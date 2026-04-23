# -*- coding: utf-8 -*-
"""
Tạo file PowerPoint từ nội dung slides.md
Chạy: py create_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──
C_PRIMARY   = RGBColor(0x02, 0x84, 0xC7)
C_DARK      = RGBColor(0x0C, 0x4A, 0x6E)
C_LIGHT     = RGBColor(0xE0, 0xF2, 0xFE)
C_ACCENT    = RGBColor(0x0E, 0xA5, 0xE9)
C_TEXT      = RGBColor(0x33, 0x41, 0x55)
C_MUTED     = RGBColor(0x64, 0x74, 0x8B)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_CODEBG    = RGBColor(0x1E, 0x29, 0x3B)
C_GREEN     = RGBColor(0x16, 0x65, 0x34)
C_RED       = RGBColor(0x99, 0x1B, 0x1B)
C_GREENBG   = RGBColor(0xDC, 0xFC, 0xE7)
C_REDBG     = RGBColor(0xFE, 0xE2, 0xE2)
C_BLUEBG    = RGBColor(0xDB, 0xEA, 0xFE)

FONT = "Calibri"
CODE_FONT = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ═══════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════

def _box(slide, l, t, w, h, fill=None, border=None, radius=True):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        l, t, w, h
    )
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh

def _tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame

def _run(para, text, sz=18, bold=False, color=C_TEXT, font=FONT):
    r = para.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return r

def _para(tf, text, sz=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, after=Pt(8), font=FONT):
    if tf.paragraphs[0].text == '' and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    _run(p, text, sz, bold, color, font)
    p.alignment = align; p.space_after = after
    return p

def _title(slide, text):
    """Top title bar with accent line"""
    _box(slide, Inches(0.6), Inches(1.1), Inches(12.1), Pt(3), fill=C_LIGHT, radius=False)
    tf = _tb(slide, Inches(0.6), Inches(0.25), Inches(12.1), Inches(0.85))
    _para(tf, text, sz=36, bold=True, color=C_PRIMARY, after=Pt(0))

def _bullets(tf, items, sz=19):
    for item in items:
        p = tf.add_paragraph()
        p.space_after = Pt(12)
        # Check for bold prefix
        if " – " in item:
            parts = item.split(" – ", 1)
            _run(p, "  \u203a  " + parts[0], sz, True, C_PRIMARY)
            _run(p, " – " + parts[1], sz, False, C_TEXT)
        elif ": " in item and len(item.split(": ")[0]) < 35:
            parts = item.split(": ", 1)
            _run(p, "  \u203a  " + parts[0] + ": ", sz, True, C_PRIMARY)
            _run(p, parts[1], sz, False, C_TEXT)
        else:
            _run(p, "  \u203a  " + item, sz, False, C_TEXT)

def _code(slide, l, t, w, h, code_text, sz=14):
    box = _box(slide, l, t, w, h, fill=C_CODEBG)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
    _para(tf, code_text, sz=sz, color=C_WHITE, font=CODE_FONT, after=Pt(0))

def _callout(slide, l, t, w, h, text, sz=16):
    _box(slide, l, t, Pt(5), h, fill=C_ACCENT, radius=False)
    box = _box(slide, l, t, w, h, fill=C_LIGHT)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    _para(tf, text, sz=sz, color=C_DARK, after=Pt(0))

def _slide_num(slide, num, total=12):
    tf = _tb(slide, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
    _para(tf, f"{num}/{total}", sz=12, color=C_MUTED, align=PP_ALIGN.RIGHT, after=Pt(0))


# ═══════════════════════════════════════
# SLIDE 1 – TRANG BÌA
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# Accent decoration
_box(s, Inches(5.6), Inches(4.2), Inches(2.1), Pt(5), fill=C_PRIMARY, radius=False)

tf = _tb(s, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.5))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
for tag in ["30 ph\u00fat", "  \u00b7  ", "SV n\u0103m 3", "  \u00b7  ", "FastAPI + Postman"]:
    _run(p, tag, 15, False, C_PRIMARY)

tf = _tb(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5))
_para(tf, "Bu\u1ed5i 8: API Testing &\nQuality Assurance", sz=56, bold=True, color=C_DARK, align=PP_ALIGN.CENTER, after=Pt(16))

tf = _tb(s, Inches(2), Inches(4.7), Inches(9.3), Inches(1))
_para(tf, "Hi\u1ec3u ph\u00e2n lo\u1ea1i test \u00b7 Vi\u1ebft test t\u1ef1 \u0111\u1ed9ng \u00b7 Ch\u1ea1y Newman & Load Test", sz=22, color=C_MUTED, align=PP_ALIGN.CENTER)

tf = _tb(s, Inches(2), Inches(5.8), Inches(9.3), Inches(0.5))
_para(tf, "Ki\u1ebfn tr\u00fac & Thi\u1ebft k\u1ebf h\u01b0\u1edbng d\u1ecbch v\u1ee5", sz=14, color=C_MUTED, align=PP_ALIGN.CENTER)

_slide_num(s, 1)


# ═══════════════════════════════════════
# SLIDE 2 – TẠI SAO PHẢI TEST API?
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "T\u1ea1i sao ph\u1ea3i Test API?")

tf = _tb(s, Inches(0.6), Inches(1.4), Inches(7.5), Inches(4.5))
_para(tf, "T\u00ecnh hu\u1ed1ng th\u1ef1c t\u1ebf", sz=26, bold=True, color=C_DARK, after=Pt(16))
_bullets(tf, [
    "Dev vi\u1ebft xong: Test tay tr\u00ean Postman \u2192 \"Ch\u1ea1y \u0111\u01b0\u1ee3c r\u1ed3i!\" \u2192 Deploy",
    "2 tu\u1ea7n sau: S\u1eeda 1 bug nh\u1ecf \u2192 API c\u0169 b\u1ecb h\u1ecfng m\u00e0 kh\u00f4ng ai bi\u1ebft (Regression Bug)",
    "Kh\u00e1ch h\u00e0ng: L\u00e0 ng\u01b0\u1eddi ph\u00e1t hi\u1ec7n l\u1ed7i \u0111\u1ea7u ti\u00ean \u2192 M\u1ea5t uy t\u00edn nghi\u00eam tr\u1ecdng",
])

_callout(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.1),
    "Nguy\u00ean t\u1eafc Shift-Left: T\u00ecm bug c\u00e0ng s\u1edbm \u1edf t\u1ea7ng API \u2192 chi ph\u00ed s\u1eeda c\u00e0ng r\u1ebb. T\u1ef1 \u0111\u1ed9ng h\u00f3a, \u0111\u1eebng ch\u1edd tester hay kh\u00e1ch h\u00e0ng.", 17)

_slide_num(s, 2)


# ═══════════════════════════════════════
# SLIDE 3 – THÁP TESTING
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "Th\u00e1p Testing (Testing Pyramid)")

# Pyramid visual
levels = [
    (Inches(4.5), Inches(1.6), Inches(4.3), Inches(0.9), "E2E / UI Test",  "Ch\u1eadm \u00b7 \u0110\u1eaft \u00b7 D\u1ec5 v\u1ee1", C_REDBG,   C_RED),
    (Inches(3.3), Inches(2.7), Inches(6.7), Inches(0.9), "Integration Test (Postman, Newman)", "Ki\u1ec3m tra li\u00ean k\u1ebft HTTP", C_BLUEBG,  RGBColor(0x1E,0x40,0xAF)),
    (Inches(2.1), Inches(3.8), Inches(9.1), Inches(0.9), "Unit Test (pytest)", "Nhanh nh\u1ea5t \u00b7 Nhi\u1ec1u nh\u1ea5t \u00b7 R\u1ebb nh\u1ea5t",  C_GREENBG, C_GREEN),
]
for l, t, w, h, title, sub, bg, fg in levels:
    box = _box(s, l, t, w, h, fill=bg)
    tf2 = box.text_frame; tf2.word_wrap = True
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, title, 18, True, fg); _run(p, "  \u2013 " + sub, 14, False, C_MUTED)

# Table
rows_data = [
    ["Lo\u1ea1i test", "Ki\u1ec3m tra c\u00e1i g\u00ec?", "T\u1ed1c \u0111\u1ed9"],
    ["Unit Test", "1 h\u00e0m / module ri\u00eang l\u1ebb", "C\u1ef1c nhanh (ms)"],
    ["Integration", "API endpoint qua HTTP", "Nhanh (gi\u00e2y)"],
    ["Performance", "Ch\u1ecbu t\u1ea3i, response time", "Ch\u1eadm (ph\u00fat)"],
]
tbl = s.shapes.add_table(4, 3, Inches(1.5), Inches(5), Inches(10.3), Inches(2.1)).table
tbl.columns[0].width = Inches(2.8); tbl.columns[1].width = Inches(5); tbl.columns[2].width = Inches(2.5)
for r in range(4):
    for c in range(3):
        cell = tbl.cell(r, c)
        cell.text = rows_data[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(15); para.font.name = FONT
            para.font.bold = (r==0); para.font.color.rgb = C_DARK if r==0 else C_TEXT

_slide_num(s, 3)


# ═══════════════════════════════════════
# SLIDE 4 – UNIT TEST VỚI PYTEST
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "Unit Test v\u1edbi pytest")

tf = _tb(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.5))
_para(tf, "\u0110\u1eb7c \u0111i\u1ec3m", sz=24, bold=True, color=C_DARK, after=Pt(12))
_bullets(tf, [
    "Test m\u1ed9t h\u00e0m duy nh\u1ea5t, bi\u1ec7t l\u1eadp",
    "Kh\u00f4ng c\u1ea7n ch\u1ea1y server, kh\u00f4ng c\u1ea7n database",
    "FastAPI h\u1ed7 tr\u1ee3 TestClient \u2013 g\u1ecdi API tr\u1ef1c ti\u1ebfp trong b\u1ed9 nh\u1edb",
    "K\u1ebft qu\u1ea3 c\u1ef1c nhanh: v\u00e0i mili-gi\u00e2y",
], sz=18)
_callout(s, Inches(0.6), Inches(5.5), Inches(5.8), Inches(0.8),
    "Ch\u1ea1y: py -m pytest test_unit.py -v", 16)

_code(s, Inches(6.7), Inches(1.4), Inches(6), Inches(4.9),
    'from fastapi.testclient import TestClient\n'
    'from main import app\n\n'
    'client = TestClient(app)  # Kh\u00f4ng c\u1ea7n uvicorn!\n\n'
    'def test_get_all_products():\n'
    '    res = client.get("/products")\n'
    '    assert res.status_code == 200\n'
    '    assert isinstance(res.json(), list)\n'
    '    assert len(res.json()) >= 1\n\n'
    'def test_product_not_found():\n'
    '    res = client.get("/products/9999")\n'
    '    assert res.status_code == 404', sz=14)

_slide_num(s, 4)


# ═══════════════════════════════════════
# SLIDE 5 – POSTMAN TEST SCRIPTS
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "Integration Test \u2013 Vi\u1ebft Test Script trong Postman")

tf = _tb(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5))
_para(tf, "Kh\u00e1c Unit Test th\u1ebf n\u00e0o?", sz=24, bold=True, color=C_DARK, after=Pt(8))
_para(tf, "G\u1ecdi API qua HTTP th\u1eadt (server \u0111ang ch\u1ea1y). Ki\u1ec3m tra routing, middleware, validation l\u00e0m vi\u1ec7c c\u00f9ng nhau.", sz=17, color=C_TEXT, after=Pt(16))
_bullets(tf, [
    'D\u00f9ng tab "Scripts > Post-response"',
    "Vi\u1ebft b\u1eb1ng JavaScript \u0111\u1ec3 Postman t\u1ef1 \u0111\u1ed9ng so s\u00e1nh k\u1ebft qu\u1ea3",
    'B\u1ecf thao t\u00e1c "nh\u00ecn b\u1eb1ng m\u1eaft" th\u1ee7 c\u00f4ng',
], sz=18)
_callout(s, Inches(0.6), Inches(5.5), Inches(5.8), Inches(1.1),
    "Hai h\u00e0m quan tr\u1ecdng nh\u1ea5t:\npm.response.to.have.status()\npm.expect().to.be.below()", 14)

_code(s, Inches(6.7), Inches(1.4), Inches(6), Inches(5.2),
    '// Ki\u1ec3m tra m\u00e3 tr\u1ea1ng th\u00e1i\n'
    'pm.test("Status code l\u00e0 200", function() {\n'
    '    pm.response.to.have.status(200);\n'
    '});\n\n'
    '// Ki\u1ec3m tra th\u1eddi gian ph\u1ea3n h\u1ed3i\n'
    'pm.test("Ph\u1ea3n h\u1ed3i d\u01b0\u1edbi 500ms", function() {\n'
    '    pm.expect(pm.response.responseTime)\n'
    '      .to.be.below(500);\n'
    '});\n\n'
    '// Ki\u1ec3m tra n\u1ed9i dung body\n'
    'pm.test("Body l\u00e0 m\u1ea3ng h\u1ee3p l\u1ec7", function() {\n'
    '    const data = pm.response.json();\n'
    '    pm.expect(data).to.be.an(\'array\');\n'
    '    pm.expect(data.length)\n'
    '      .to.be.greaterThan(0);\n'
    '});', sz=13)

_slide_num(s, 5)


# ═══════════════════════════════════════
# SLIDE 6 – COLLECTION VARIABLES
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "Collection Variables \u2013 Chu\u1ed7i Request li\u00ean k\u1ebft")

tf = _tb(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1))
_para(tf, "POST t\u1ea1o s\u1ea3n ph\u1ea9m \u2192 Server sinh ID m\u1edbi \u2192 L\u00e0m sao d\u00f9ng \u0111\u00fang ID \u0111\u00f3 cho GET, PUT, DELETE?\nD\u00f9ng Collection Variables truy\u1ec1n d\u1eef li\u1ec7u gi\u1eefa c\u00e1c request.",
    sz=19, color=C_TEXT, align=PP_ALIGN.CENTER, after=Pt(0))

# Flow steps
steps = [
    ("1. POST /products", "T\u1ea1o m\u1edbi s\u1ea3n ph\u1ea9m\nL\u01b0u created_id = 4"),
    ("2. GET /products/{id}", "T\u1ef1 \u0111i\u1ec1n ID=4\nKi\u1ec3m tra t\u1ed3n t\u1ea1i"),
    ("3. PUT /products/{id}", "C\u1eadp nh\u1eadt gi\u00e1\ns\u1ea3n ph\u1ea9m \u0111\u00f3"),
    ("4. DELETE /{id}", "X\u00f3a s\u1ea3n ph\u1ea9m\nGET l\u1ea1i \u2192 404 \u2713"),
]
for i, (title, desc) in enumerate(steps):
    x = Inches(0.4) + Inches(3.15) * i
    box = _box(s, x, Inches(2.8), Inches(2.8), Inches(2), fill=C_WHITE, border=C_LIGHT)
    tf2 = box.text_frame; tf2.word_wrap = True
    _para(tf2, title, sz=16, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER, after=Pt(6))
    _para(tf2, desc, sz=14, color=C_MUTED, align=PP_ALIGN.CENTER, after=Pt(0))
    if i < 3:
        arr = _box(s, x + Inches(2.85), Inches(3.55), Inches(0.35), Inches(0.35), fill=C_PRIMARY)
        # Arrow label
        tf3 = _tb(s, x + Inches(2.82), Inches(3.45), Inches(0.4), Inches(0.5))
        _para(tf3, "\u2192", sz=22, color=C_PRIMARY, align=PP_ALIGN.CENTER, after=Pt(0))

_code(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.8),
    '// Trong Tests c\u1ee7a POST \u2013 L\u01b0u ID v\u1eeba t\u1ea1o:\n'
    'const data = pm.response.json();\n'
    'pm.collectionVariables.set("created_id", data.id);\n\n'
    '// Trong URL c\u1ee7a GET/PUT/DELETE \u2013 T\u1ef1 \u0111i\u1ec1n:\n'
    '{{base_url}}/products/{{created_id}}', sz=14)

_slide_num(s, 6)


# ═══════════════════════════════════════
# SLIDE 7 – NEWMAN
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "Newman \u2013 \u0110\u01b0a Postman v\u00e0o Terminal & CI/CD")

tf = _tb(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(4))
_para(tf, "Newman l\u00e0 g\u00ec?", sz=24, bold=True, color=C_DARK, after=Pt(12))
_bullets(tf, [
    "Phi\u00ean b\u1ea3n d\u00f2ng l\u1ec7nh (CLI) c\u1ee7a Postman Runner",
    "C\u00f9ng file .json, c\u00f9ng test script \u2013 ch\u1ea1y trong Terminal",
    "T\u00edch h\u1ee3p GitHub Actions / GitLab CI: Ch\u1eb7n code l\u1ed7i tr\u01b0\u1edbc khi deploy",
], sz=18)
_callout(s, Inches(0.6), Inches(5.2), Inches(5.8), Inches(1.1),
    "C\u00e0i: npm install -g newman\nCh\u1ea1y: newman run postman_collection.json", 15)

_code(s, Inches(6.7), Inches(1.4), Inches(6), Inches(5.2),
    '$ newman run postman_collection.json\n\n'
    '\u2192 1. GET /products - L\u1ea5y danh s\u00e1ch\n'
    '  GET localhost:8000/products [200 OK, 45ms]\n'
    '  \u2713 Status code = 200\n'
    '  \u2713 Response time < 500ms\n'
    '  \u2713 Body tr\u1ea3 v\u1ec1 l\u00e0 m\u1ea3ng\n\n'
    '\u2192 2. POST /products - T\u1ea1o m\u1edbi\n'
    '  POST localhost:8000/products [201 Created]\n'
    '  \u2713 Server t\u1ef1 sinh ID\n\n'
    '\u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510\n'
    '\u2502 assertions     \u2502   14  \u2502       \u2502\n'
    '\u2502 failed         \u2502    0  \u2502       \u2502\n'
    '\u2502 total duration \u2502 210ms \u2502       \u2502\n'
    '\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518', sz=13)

_slide_num(s, 7)


# ═══════════════════════════════════════
# SLIDE 8 – PERFORMANCE TESTING
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "Performance Testing \u2013 3 Ch\u1ec9 s\u1ed1 sinh t\u1ed3n")

tf = _tb(s, Inches(1), Inches(1.4), Inches(11.3), Inches(0.7))
_para(tf, "API ch\u1ea1y \u0111\u00fang v\u1edbi 1 user \u2013 nh\u01b0ng n\u1ebfu 1000 ng\u01b0\u1eddi g\u1ecdi c\u00f9ng l\u00fac th\u00ec server c\u00f3 s\u1eadp kh\u00f4ng?",
    sz=21, color=C_TEXT, align=PP_ALIGN.CENTER, after=Pt(0))

metrics = [
    ("Response Time", "Th\u1eddi gian t\u1eeb g\u1eedi request \u0111\u1ebfn khi nh\u1eadn response.\n< 200ms (t\u1ed1t) \u00b7 > 3s (th\u1ea3m ho\u1ea1)", C_ACCENT),
    ("Throughput (RPS)", "S\u1ed1 request server x\u1eed l\u00fd trong 1 gi\u00e2y.\nC\u00e0ng cao c\u00e0ng t\u1ed1t.", RGBColor(0x22,0xC5,0x5E)),
    ("Error Rate", "T\u1ef7 l\u1ec7 request b\u1ecb l\u1ed7i (5xx).\n< 1% (t\u1ed1t) \u00b7 > 5% (nguy k\u1ecbch)", RGBColor(0xEF,0x44,0x44)),
]
for i, (title, desc, accent) in enumerate(metrics):
    x = Inches(0.5) + Inches(4.2) * i
    _box(s, x, Inches(2.3), Inches(3.8), Pt(6), fill=accent, radius=False)
    box = _box(s, x, Inches(2.3), Inches(3.8), Inches(3.5), fill=C_WHITE, border=C_LIGHT)
    tf2 = box.text_frame; tf2.word_wrap = True
    tf2.margin_top = Inches(0.3)
    _para(tf2, title, sz=24, bold=True, color=C_DARK, align=PP_ALIGN.CENTER, after=Pt(14))
    _para(tf2, desc, sz=16, color=C_MUTED, align=PP_ALIGN.CENTER, after=Pt(0))

_callout(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
    "Load Test: N user \u0111\u1ed3ng th\u1eddi  \u00b7  Stress Test: T\u0103ng d\u1ea7n \u0111\u1ebfn s\u1eadp  \u00b7  Spike Test: 1000 user \u0111\u1ed9t ng\u1ed9t (flash sale)", 16)

_slide_num(s, 8)


# ═══════════════════════════════════════
# SLIDE 9 – LOCUST
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "Locust \u2013 Load Testing b\u1eb1ng Python")

tf = _tb(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.5))
_para(tf, "T\u1ea1i sao ch\u1ecdn Locust?", sz=24, bold=True, color=C_DARK, after=Pt(12))
_bullets(tf, [
    "Vi\u1ebft k\u1ecbch b\u1ea3n b\u1eb1ng Python (quen thu\u1ed9c!)",
    "@task(n): tr\u1ecdng s\u1ed1 h\u00e0nh vi user",
    "wait_time: m\u00f4 ph\u1ecfng user th\u1eadt c\u00f3 ngh\u1ec9",
    "Web UI real-time \u0111\u1ed3 th\u1ecb c\u1ef1c \u0111\u1eb9p",
], sz=18)
_callout(s, Inches(0.6), Inches(5.3), Inches(5.8), Inches(1.2),
    "Ch\u1ea1y: locust\nM\u1edf: http://localhost:8089\n\u0110i\u1ec1n: 50 users \u00b7 Host: http://localhost:8000", 14)

_code(s, Inches(6.7), Inches(1.4), Inches(6), Inches(5.2),
    'from locust import HttpUser, task, between\n\n'
    'class ProductAPIUser(HttpUser):\n'
    '    wait_time = between(0.5, 1.5)\n\n'
    '    @task(5)  # G\u1ecdi nhi\u1ec1u nh\u1ea5t\n'
    '    def browse_products(self):\n'
    '        self.client.get("/products")\n\n'
    '    @task(1)  # G\u1ecdi \u00edt h\u01a1n\n'
    '    def create_product(self):\n'
    '        self.client.post("/products",\n'
    '            json={\n'
    '                "name": "Test",\n'
    '                "price": 100000\n'
    '            }\n'
    '        )', sz=14)

_slide_num(s, 9)


# ═══════════════════════════════════════
# SLIDE 10 – BẢNG SO SÁNH
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "T\u1ed5ng k\u1ebft \u2013 Ch\u1ecdn c\u00f4ng c\u1ee5 n\u00e0o?")

rows_data = [
    ["C\u00f4ng c\u1ee5", "Lo\u1ea1i test", "Ng\u00f4n ng\u1eef", "C\u1ea7n server?", "CI/CD"],
    ["pytest",   "Unit Test",    "Python",     "\u2717 Kh\u00f4ng", "\u2713 R\u1ea5t d\u1ec5"],
    ["Postman",  "Integration",  "JavaScript", "\u2713 C\u00f3",     "\u2717 GUI only"],
    ["Newman",   "Integration",  "JavaScript", "\u2713 C\u00f3",     "\u2713 Xu\u1ea5t s\u1eafc"],
    ["Locust",   "Performance",  "Python",     "\u2713 C\u00f3",     "\u2713 T\u1ed1t"],
]
tbl = s.shapes.add_table(5, 5, Inches(0.6), Inches(1.5), Inches(12.1), Inches(3.2)).table
for c in range(5):
    tbl.columns[c].width = Inches(2.42)
for r in range(5):
    for c in range(5):
        cell = tbl.cell(r, c)
        cell.text = rows_data[r][c]
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(16); para.font.name = FONT
            para.font.bold = (r==0 or c==0)
            para.font.color.rgb = C_DARK if r==0 else (C_GREEN if "\u2713" in rows_data[r][c] else (C_RED if "\u2717" in rows_data[r][c] else C_TEXT))

# Pipeline
tf = _tb(s, Inches(0.6), Inches(5), Inches(12.1), Inches(0.5))
_para(tf, "Quy tr\u00ecnh l\u00fd t\u01b0\u1edfng trong d\u1ef1 \u00e1n th\u1ef1c t\u1ebf:", sz=20, bold=True, color=C_DARK, align=PP_ALIGN.CENTER, after=Pt(8))

pipeline = "Dev Code  \u2192  pytest  \u2192  Push Git  \u2192  CI ch\u1ea1y Newman  \u2192  Staging  \u2192  Locust  \u2192  Production \u2713"
box = _box(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.9), fill=C_LIGHT)
tf2 = box.text_frame; tf2.word_wrap = True
_para(tf2, pipeline, sz=18, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER, after=Pt(0))

_slide_num(s, 10)


# ═══════════════════════════════════════
# SLIDE 11 – CÂU HỎI THƯỜNG GẶP
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
_title(s, "C\u00e2u h\u1ecfi th\u01b0\u1eddng g\u1eb7p")

qas = [
    ("Unit Test v\u1edbi Integration Test kh\u00e1c nhau ch\u1ed7 n\u00e0o?",
     "Unit Test d\u00f9ng TestClient x\u1eed l\u00fd trong b\u1ed9 nh\u1edb (kh\u00f4ng qua HTTP, nhanh g\u1ea5p 100x). Integration Test g\u1ecdi qua HTTP th\u1eadt \u2192 b\u1eaft \u0111\u01b0\u1ee3c l\u1ed7i routing, CORS, middleware."),
    ("Locust vi\u1ebft b\u1eb1ng Python, v\u1eady test API Java/Go \u0111\u01b0\u1ee3c kh\u00f4ng?",
     "Locust ch\u1ec9 g\u1eedi HTTP request \u0111\u1ebfn URL. N\u00f3 kh\u00f4ng quan t\u00e2m backend vi\u1ebft b\u1eb1ng ng\u00f4n ng\u1eef g\u00ec. Test \u0111\u01b0\u1ee3c b\u1ea5t k\u1ef3 API n\u00e0o."),
    ("Sao kh\u00f4ng d\u00f9ng Postman lu\u00f4n, c\u1ea7n chi Newman?",
     "Postman c\u1ea7n ng\u01b0\u1eddi m\u1edf app b\u1ea5m Run. Newman ch\u1ea1y CLI \u2192 t\u00edch h\u1ee3p CI/CD. M\u1ed7i l\u1ea7n push code, Newman t\u1ef1 ch\u1ea1y kh\u00f4ng c\u1ea7n con ng\u01b0\u1eddi."),
    ("collectionVariables kh\u00e1c environment th\u1ebf n\u00e0o?",
     "collectionVariables g\u1eafn v\u1edbi file JSON, d\u1ec5 chia s\u1ebb. environment ch\u1ee9a config nh\u1ea1y c\u1ea3m (API keys), t\u00e1ch bi\u1ec7t \u0111\u1ec3 b\u1ea3o m\u1eadt."),
]
y = Inches(1.5)
for q, a in qas:
    # Question
    tf = _tb(s, Inches(0.8), y, Inches(11.7), Inches(0.4))
    _para(tf, "Q: " + q, sz=17, bold=True, color=C_PRIMARY, after=Pt(2))
    # Answer with left border
    _box(s, Inches(0.8), y + Inches(0.45), Pt(4), Inches(0.65), fill=C_LIGHT, radius=False)
    tf = _tb(s, Inches(1.1), y + Inches(0.45), Inches(11.4), Inches(0.65))
    _para(tf, a, sz=15, color=C_TEXT, after=Pt(0))
    y += Inches(1.35)

_slide_num(s, 11)


# ═══════════════════════════════════════
# SLIDE 12 – Q&A
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)

_box(s, Inches(5.6), Inches(4.5), Inches(2.1), Pt(5), fill=C_PRIMARY, radius=False)

tf = _tb(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4))
_para(tf, "?", sz=60, align=PP_ALIGN.CENTER, after=Pt(8), color=C_PRIMARY)
_para(tf, "Q & A", sz=72, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER, after=Pt(20))
_para(tf, "C\u1ea3m \u01a1n c\u00e1c b\u1ea1n \u0111\u00e3 theo d\u00f5i.\nM\u1ecdi th\u1eafc m\u1eafc xin m\u1eddi \u0111\u1eb7t c\u00e2u h\u1ecfi!", sz=24, color=C_MUTED, align=PP_ALIGN.CENTER)

_slide_num(s, 12)


# ═══════════════════════════════════════
# SAVE
# ═══════════════════════════════════════
out = "Slides_Buoi_8.pptx"
prs.save(out)
import sys
sys.stdout.buffer.write(f"OK -> {out}\n".encode("utf-8"))
